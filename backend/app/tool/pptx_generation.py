import os
import re
from typing import List, Dict, Any, Optional

from app.tool.base import BaseTool, ToolResult
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.shapes import MSO_SHAPE


class PptxGenerationTool(BaseTool):
    name: str = "create_presentation"
    description: str = "Creates an advanced PowerPoint presentation. supports rich text, charts, tables, and shapes. ALWAYS use the specific parameters for charts/tables, do not put them in text content."
    parameters: dict = {
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "The filename for the presentation. Should end with .pptx.",
            },
            "slides": {
                "type": "array",
                "description": "A list of slides to add to the presentation.",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string", "description": "The title of the slide."},
                        "layout": {
                            "type": "string",
                            "description": "Layout type: 'title_slide', 'title_and_content', 'two_content', 'comparison', 'picture_with_caption', 'blank'. Defaults to 'title_and_content'.",
                            "default": "title_and_content",
                        },
                        "content": {
                            "type": "string",
                            "description": "The text content. Supports **bold**, *italic*, and __underline__ formatting. Do NOT put tables or charts here. Use the dedicated 'tables' and 'charts' parameters.",
                        },
                        "font_name": {"type": "string", "description": "Font name (e.g., 'Arial')."},
                        "font_size": {"type": "integer", "description": "Font size in points."},
                        "font_color": {"type": "string", "description": "Hex code for font color."},
                        "background_color": {"type": "string", "description": "Hex code for background color."},
                        "charts": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "type": {"type": "string", "enum": ["bar", "line", "pie"]},
                                    "title": {"type": "string"},
                                    "data": {"type": "object", "description": "Dictionary of category -> value (for pie) or category -> list of values (for bar/line)"},
                                    "series_names": {"type": "array", "items": {"type": "string"}, "description": "List of series names for bar/line charts"}
                                }
                            }
                        },
                        "tables": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "rows": {"type": "integer"},
                                    "cols": {"type": "integer"},
                                    "data": {"type": "array", "items": {"type": "array", "items": {"type": "string"}}}
                                }
                            }
                        },
                        "shapes": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "type": {"type": "string", "enum": ["rectangle", "circle", "arrow", "triangle"]},
                                    "text": {"type": "string"},
                                    "color": {"type": "string", "description": "Hex color code"}
                                }
                            }
                        }
                    },
                    "required": ["title"],
                },
            },
        },
        "required": ["filename", "slides"],
    }

    def _apply_rich_text(self, paragraph, text: str, font_name: str = "Arial", font_size: int = 18, font_color: str = None):
        """Parses markdown-like syntax and applies formatting to runs."""
        # Regex for **bold**, *italic*, __underline__
        # This is a simple parser and might not handle nested tags perfectly
        tokens = re.split(r'(\*\*.*?\*\*|\*.*?\*|__.*?__)', text)
        
        for token in tokens:
            if not token:
                continue
                
            run = paragraph.add_run()
            
            # Check formatting
            if token.startswith('**') and token.endswith('**'):
                run.text = token[2:-2]
                run.font.bold = True
            elif token.startswith('*') and token.endswith('*'):
                run.text = token[1:-1]
                run.font.italic = True
            elif token.startswith('__') and token.endswith('__'):
                run.text = token[2:-2]
                run.font.underline = True
            else:
                run.text = token
                
            # Apply common font properties
            run.font.name = font_name
            run.font.size = Pt(font_size)
            if font_color and re.match(r"^[0-9a-fA-F]{6}$", font_color):
                run.font.color.rgb = RGBColor.from_string(font_color)

    def _add_chart(self, slide, chart_data: Dict[str, Any]):
        chart_type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE
        }
        
        c_type = chart_type_map.get(chart_data.get("type", "bar"))
        data_obj = CategoryChartData() if c_type != XL_CHART_TYPE.PIE else ChartData()
        
        raw_data = chart_data.get("data", {})
        
        if c_type == XL_CHART_TYPE.PIE:
            data_obj.categories = list(raw_data.keys())
            data_obj.add_series(chart_data.get("title", "Series"), list(raw_data.values()))
        else:
            # Bar/Line: Expecting data as {category: [val1, val2]} or simple {category: val}
            categories = list(raw_data.keys())
            data_obj.categories = categories
            
            series_names = chart_data.get("series_names", ["Series 1"])
            
            # Normalize data to list of lists
            first_val = list(raw_data.values())[0]
            if not isinstance(first_val, list):
                # Single series
                data_obj.add_series(series_names[0], list(raw_data.values()))
            else:
                # Multiple series
                num_series = len(first_val)
                for i in range(num_series):
                    name = series_names[i] if i < len(series_names) else f"Series {i+1}"
                    values = [raw_data[cat][i] for cat in categories]
                    data_obj.add_series(name, values)

        # Add chart to slide (bottom right quadrant approx)
        slide.shapes.add_chart(
            c_type, 
            Inches(5), Inches(2), Inches(4.5), Inches(3), 
            data_obj
        )

    def _add_table(self, slide, table_data: Dict[str, Any]):
        rows = table_data.get("rows", 2)
        cols = table_data.get("cols", 2)
        data = table_data.get("data", [])
        
        # Add table (bottom left quadrant approx)
        shape = slide.shapes.add_table(
            rows, cols, Inches(0.5), Inches(2), Inches(4), Inches(0.8 * rows)
        )
        table = shape.table
        
        for r in range(min(rows, len(data))):
            for c in range(min(cols, len(data[r]))):
                table.cell(r, c).text = str(data[r][c])

    def _add_shape(self, slide, shape_data: Dict[str, Any]):
        shape_type_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "circle": MSO_SHAPE.OVAL,
            "arrow": MSO_SHAPE.RIGHT_ARROW,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE
        }
        
        s_type = shape_type_map.get(shape_data.get("type", "rectangle"))
        
        # Add shape (random-ish position or fixed for now)
        shape = slide.shapes.add_shape(
            s_type, Inches(7), Inches(5), Inches(2), Inches(1)
        )
        
        if shape_data.get("text"):
            shape.text = shape_data.get("text")
            
        if shape_data.get("color"):
            color = shape_data.get("color")
            if re.match(r"^[0-9a-fA-F]{6}$", color):
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(color)

    async def execute(self, filename: str, slides: list) -> ToolResult:
        if not filename.endswith(".pptx"):
            return ToolResult(error="Filename must end with .pptx")

        try:
            prs = Presentation()
            
            # Expanded layout map
            layout_map = {
                "title_slide": 0,
                "title_and_content": 1,
                "two_content": 3,
                "comparison": 4,
                "title_only": 5,
                "blank": 6,
                "picture_with_caption": 8
            }

            for slide_data in slides:
                layout_name = slide_data.get("layout", "title_and_content").lower()
                layout_index = layout_map.get(layout_name, 1)
                
                # Handle case where layout index might be out of bounds for some templates
                # Default to 1 (Title and Content) if not found
                if layout_index >= len(prs.slide_layouts):
                    layout_index = 1
                    
                slide_layout = prs.slide_layouts[layout_index]
                slide = prs.slides.add_slide(slide_layout)

                # 1. Set Title
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.get("title", "")

                # 2. Set Content (Text)
                # For standard layouts, placeholder[1] is usually the content body
                content_text = slide_data.get("content", "")
                if content_text:
                    # Try to find the body placeholder
                    body_shape = None
                    for shape in slide.placeholders:
                        if shape.placeholder_format.idx == 1:
                            body_shape = shape
                            break
                    
                    if body_shape and hasattr(body_shape, "text_frame"):
                        tf = body_shape.text_frame
                        tf.clear()
                        
                        font_name = slide_data.get("font_name", "Arial")
                        font_size = slide_data.get("font_size", 18)
                        font_color = slide_data.get("font_color")

                        lines = content_text.split("\n")
                        for i, line in enumerate(lines):
                            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                            p.level = 0 if not line.startswith("  ") else 1 # Simple indentation
                            clean_line = line.strip()
                            self._apply_rich_text(p, clean_line, font_name, font_size, font_color)
                            
                # 3. Add Charts
                if "charts" in slide_data:
                    for chart_data in slide_data["charts"]:
                        self._add_chart(slide, chart_data)
                        
                # 4. Add Tables
                if "tables" in slide_data:
                    for table_data in slide_data["tables"]:
                        self._add_table(slide, table_data)
                        
                # 5. Add Shapes
                if "shapes" in slide_data:
                    for shape_data in slide_data["shapes"]:
                        self._add_shape(slide, shape_data)

                # 6. Background Color
                bg_color_hex = slide_data.get("background_color")
                if bg_color_hex and re.match(r"^[0-9a-fA-F]{6}$", bg_color_hex):
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor.from_string(bg_color_hex)

            workspace_dir = os.path.join(os.getcwd(), "workspace")
            if not os.path.exists(workspace_dir):
                os.makedirs(workspace_dir)

            filepath = os.path.join(workspace_dir, filename)
            prs.save(filepath)

            return ToolResult(output=f"Presentation successfully created at {filepath}")

        except Exception as e:
            return ToolResult(error=f"Failed to create presentation: {str(e)}")
